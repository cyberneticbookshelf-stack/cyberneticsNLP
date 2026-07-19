#!/usr/bin/env python3
"""patch_deck_566.py — migrate the talk deck from the 26 April taxonomy to the
19 July 2026 566-book re-canonicalisation (run_20260719_k9_s5, class
88c44bece9a5a875).

This is the successor to patch_deck.py (which was a one-shot v2->v3 April
migration and is not re-runnable: its input v2.pptx is gone and its find/replace
targets are pre-April strings). This script does targeted, per-slide,
paragraph-level find/replace.

Reads:  presentation/CyberneticsNLP_Talk_v4.pptx   (April taxonomy)
Writes: presentation/CyberneticsNLP_Talk_v4_566.pptx   (non-destructive)

What it changes (27 edits, all verified applied 19 Jul 2026):
  - slide 4  : corpus headline 733·541 -> 739·566; footnote de-duplicated
               (the "541 analysed — [2133] excluded" sentence was repeated 5x)
               and rewritten to "566 analysed (OCR-corrupt books excluded)"
               ([2133] id dropped — reassigned by the July reconstruction).
  - slide 5,24: narrative "733 books" -> "566 books".
  - slide 10 : API-cost line "541 books" -> "566 books".
  - slide 11 : stability 0.348 / 5-of-9 / T7-unstable -> 0.365 / 6-of-9 /
               T1-unstable (0.145). Bands are 09c thresholds (>=0.30 stable);
               log_pipeline_run.py counts 4 stable — KI-11, an open inconsistency.
  - slide 13 : the 9 topic names AND their 7 shown top-word lines, positionally
               (T1..T9). Positions permuted vs April, so names AND keywords move.
  - slide 14 : the four era -> dominant-topic lines, remapped from the new run's
               era x dominant-topic distribution:
                 Foundational (<=1969):  T8 Control&Feedback(21) · T7 Management(5)
                 Second-Order (1970s):   T5 Social Systems · T8 Control&Feedback
                 Social Scale (80s-90s): T5 Social Systems(17) · T7 Management(16)
                 Diffusion (2000s-25):   T5 Social Systems(91) · T9 Digital Arts(81)
               EDITORIAL NOTE (Paul Wong, 19 Jul): the Second-Order Wave line
               features T5 (Social Systems and Second-Order Constructivism)
               because it matches the era's intellectual identity (Varela /
               autopoiesis), even though the raw 1970s book-count leaders were
               T8/T7/T6. Kept deliberately; not a data-lead choice.

Names T1..T9 (finalised 19 Jul, provisional single-rater): History of
Information Age and Cybernetics · Extensions and Exploration of Cybernetics ·
Biological and Ecological Regulation: Homeostasis & Allostasis · Cybernetics of
Self · Social Systems and Second-Order Constructivism · Foundations of
Cybernetics · Management and Organisational Cybernetics · Control and Feedback
Systems · Digital Arts, Architecture, Design and Posthumanism.

Usage:  python3 presentation/patch_deck_566.py   (run from repo root)
Requires: python-pptx.
"""
from pptx import Presentation

SRC = 'presentation/CyberneticsNLP_Talk_v4.pptx'
OUT = 'presentation/CyberneticsNLP_Talk_v4_566.pptx'


def replace_on_slide(pres, idx, old, new):
    """Replace old->new in any paragraph on slide idx. Paragraph-level:
    collapses the match to run 0 (preserving that run's formatting)."""
    hit = False
    for sh in pres.slides[idx].shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if old in full:
                newfull = full.replace(old, new)
                if para.runs:
                    para.runs[0].text = newfull
                    for r in para.runs[1:]:
                        r.text = ""
                hit = True
    return hit


def set_paragraph(pres, idx, marker, newtext):
    """Rewrite the whole paragraph containing `marker` to `newtext`."""
    for sh in pres.slides[idx].shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            full = "".join(r.text for r in para.runs)
            if marker in full:
                if para.runs:
                    para.runs[0].text = newtext
                    for r in para.runs[1:]:
                        r.text = ""
                return True
    return False


# (slide_idx, old, new)
REPL = [
    # slide 4 — corpus headline
    (4, "733 · 541", "739 · 566"),
    # slide 10 — API cost line book count
    (10, "541 books", "566 books"),
    # slide 11 — stability figures (two phrasings)
    (11, "mean stability 0.348, 5/9 topics stable (T7 unstable, 0.045)",
         "mean stability 0.365, 6/9 topics stable (T1 unstable, 0.145)"),
    (11, "5-seed stability complete (mean=0.348, 5/9 stable; T7 unstable)",
         "5-seed stability complete (mean=0.365, 6/9 stable; T1 unstable)"),
    # slide 13 — nine topic names (positional T1..T9)
    (13, "History and Historiography of Cybernetics", "History of Information Age and Cybernetics"),
    (13, "Techno-political Complexes", "Extensions and Exploration of Cybernetics"),
    (13, "Engineering Control", "Biological and Ecological Regulation: Homeostasis & Allostasis"),
    (13, "Social and Organisational Cybernetics", "Cybernetics of Self"),
    (13, "Formal Foundations of Cybernetics", "Social Systems and Second-Order Constructivism"),
    (13, "Reinventing Selves and Others, Past and Future", "Foundations of Cybernetics"),
    (13, "Psychological and Behavioural Regulation and Control", "Management and Organisational Cybernetics"),
    (13, "T8 Biological and Neural Cybernetics", "T8 Control and Feedback Systems"),
    (13, "T9 Extensions of Cybernetics", "T9 Digital Arts, Architecture, Design and Posthumanism"),
    # slide 14 — era -> dominant-topic remap (data-driven; Second-Order = editorial, see header)
    (14, "Dominant: T3 Engineering Control · T5 Formal Foundations of Cybernetics",
         "Dominant: T8 Control and Feedback Systems · T7 Management and Organisational Cybernetics"),
    (14, "Dominant: T9 Extensions of Cybernetics · T8 Biological and Neural Cybernetics",
         "Dominant: T5 Social Systems and Second-Order Constructivism · T8 Control and Feedback Systems"),
    (14, "Dominant: T4 Social and Organisational Cybernetics · T3 Engineering Control",
         "Dominant: T5 Social Systems and Second-Order Constructivism · T7 Management and Organisational Cybernetics"),
    (14, "Dominant: T9 Extensions of Cybernetics · T6 Reinventing Selves and Others, Past and Future",
         "Dominant: T5 Social Systems and Second-Order Constructivism · T9 Digital Arts, Architecture, Design and Posthumanism"),
    # slide 13 — per-topic top-word lines (positional, new run)
    (13, "machine · computer · wiener · cybernetic · brain · intelligence",
         "machine · computer · wiener · century · technology · american"),
    (13, "technology · cybernetic · project · architecture · machine · computer",
         "voice · sound · music · qian · chinese · china"),
    (13, "variable · input · output · equation · rate · feedback",
         "brain · cell · animal · organism · evolution · energy"),
    (13, "social · organization · decision · society · communication · complexity",
         "person · feel · bateson · child · family · behavior"),
    (13, "language · machine · object · probability · message · entropy",
         "social · communication · language · meaning · object · distinction"),
    (13, "family · person · tell · feel · child · therapy",
         "define · entropy · probability · theorem · equation · shall"),
    (13, "behavior · brain · emotion · social · person · mental",
         "organization · management · social · decision · variety · cybernetic"),
    # narrative book count (analysed corpus)
    (5, "733 books", "566 books"),
    (24, "733 books, 70 years", "566 books, 70 years"),
]

# (slide_idx, marker_substring, new_full_paragraph_text)
PARA_SET = [
    (4, "GST dates from the 1950s",
     "* GST dates from the 1950s; corpus coverage of that period may be incomplete.  "
     "Findings are provisional — the collection is not exhaustive.  "
     "566 analysed (OCR-corrupt books excluded)."),
]


def main():
    pres = Presentation(SRC)
    misses = []
    for idx, old, new in REPL:
        if replace_on_slide(pres, idx, old, new):
            print(f"  ok    slide {idx}: {old[:48]!r} -> {new[:40]!r}")
        else:
            misses.append((idx, old))
            print(f"  MISS  slide {idx}: {old[:64]!r}")
    for idx, marker, newtext in PARA_SET:
        if set_paragraph(pres, idx, marker, newtext):
            print(f"  ok    slide {idx}: paragraph rewrite ({marker[:30]!r})")
        else:
            misses.append((idx, marker))
            print(f"  MISS  slide {idx}: paragraph marker {marker[:30]!r}")
    pres.save(OUT)
    total = len(REPL) + len(PARA_SET)
    print(f"\nSaved {OUT}  ({total - len(misses)}/{total} edits applied)")
    if misses:
        print(f"WARNING: {len(misses)} strings not found — review those slides.")


if __name__ == '__main__':
    main()
