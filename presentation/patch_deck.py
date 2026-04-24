#!/usr/bin/env python3
"""
patch_deck.py — apply v0.5.1 canonical facts to CyberneticsNLP_Talk_v2.pptx.

Usage (on the NLP machine, from project root):
    python3 presentation/patch_deck.py

Reads:  presentation/CyberneticsNLP_Talk_v2.pptx
        json/nlp_results.json       (topic names + top words)
        json/entity_network.json    (node/edge counts — optional)
        src/15_entity_classify.py   (for MANUAL_CORRECTIONS count — optional)

Writes: presentation/CyberneticsNLP_Talk_v3.pptx
"""

import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
except ImportError:
    sys.exit("python-pptx missing. Run: conda install main::python-pptx")

ROOT = Path(__file__).resolve().parent.parent
# Once v3 exists, edit it in place so manual tweaks (font sizes etc.) persist.
# If v3 is missing (first-time run), bootstrap from v2.
V2 = ROOT / 'presentation' / 'CyberneticsNLP_Talk_v2.pptx'
V3 = ROOT / 'presentation' / 'CyberneticsNLP_Talk_v3.pptx'
SRC = V3 if V3.exists() else V2
DST = V3

# ── Canonical topic order (must match the current canonical run) ─────────────
# If json/nlp_results.json provides topic_names keyed differently, we map into
# this canonical sequence. T1 → index 0.
CANONICAL_NAMES = [
    'History and Biography of Cybernetics',
    'Cybernetics of Psychology',
    'Extensions of Cybernetics',
    'Cybernetic Management Theory',
    'Biological Systems Cybernetics',
    'Formal Foundations of Cybernetics',
    'Cross-Domain Applications of Cybernetics',
    'Cybernetics of Posthumanism',
    'Cultural Applications of Cybernetics',
]
N_WORDS = 6  # top words per topic shown on slide

# ── Helpers ───────────────────────────────────────────────────────────────────

def replace_in_runs(shape, old, new):
    """Replace `old` with `new` within a text frame's runs (preserves formatting).
    Returns True on the first hit."""
    if not shape.has_text_frame:
        return False
    for para in shape.text_frame.paragraphs:
        # Collapse paragraph run text to spot the substring even if it spans runs.
        full = "".join(r.text for r in para.runs)
        if old not in full:
            continue
        # Simple case: single run contains the substring.
        for run in para.runs:
            if old in run.text:
                run.text = run.text.replace(old, new)
                return True
        # Spans multiple runs: rewrite first run, clear the rest.
        para.runs[0].text = full.replace(old, new)
        for r in para.runs[1:]:
            r.text = ""
        return True
    return False


def set_shape_text(shape, new_text):
    """Replace a text frame's text preserving the first run's formatting."""
    if not shape.has_text_frame:
        return False
    tf = shape.text_frame
    # Keep first paragraph, wipe others.
    first_para = tf.paragraphs[0]
    if first_para.runs:
        first_para.runs[0].text = new_text
        for r in first_para.runs[1:]:
            r.text = ""
    else:
        first_para.add_run().text = new_text
    # Remove extra paragraphs via XML (python-pptx has no direct delete).
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    return True


def walk_and_replace(pres, slide_idx, old, new):
    """Hunt across all shapes of a slide for the substring and replace."""
    slide = pres.slides[slide_idx]
    for shape in slide.shapes:
        if replace_in_runs(shape, old, new):
            return True
    return False


# ── Load live facts ───────────────────────────────────────────────────────────

def load_topic_data():
    """Return list of (name, [top words]) for T1..T9 in canonical order.
    Falls back to hardcoded canonical names + empty words if nlp_results.json
    cannot be parsed into the expected shape."""
    path = ROOT / 'json' / 'nlp_results.json'
    words_per_topic = [[] for _ in range(9)]
    names = list(CANONICAL_NAMES)

    try:
        with open(path) as f:
            nlp = json.load(f)
    except FileNotFoundError:
        print(f"  [warn] {path} not found — using hardcoded names, empty words")
        return list(zip(names, words_per_topic))

    # Topic names — try common shapes.
    tn = nlp.get('topic_names')
    if isinstance(tn, list) and len(tn) >= 9:
        names = [str(tn[i]) for i in range(9)]
    elif isinstance(tn, dict):
        # Keys may be 'T1'..'T9' or '0'..'8' or '1'..'9'.
        ordered = []
        for k in range(1, 10):
            for candidate in (f'T{k}', str(k), str(k - 1)):
                if candidate in tn:
                    ordered.append(str(tn[candidate]))
                    break
            else:
                ordered.append(CANONICAL_NAMES[k - 1])
        names = ordered

    # Top words — try several key names.
    for key in ('top_words', 'topic_words', 'topics', 'words_per_topic'):
        tw = nlp.get(key)
        if tw is None:
            continue
        if isinstance(tw, list) and len(tw) >= 9:
            for i in range(9):
                item = tw[i]
                if isinstance(item, list):
                    words_per_topic[i] = [str(w) for w in item[:N_WORDS]]
                elif isinstance(item, dict):
                    # e.g. {'words': [...], ...}
                    ws = item.get('words') or item.get('top_words') or []
                    words_per_topic[i] = [str(w) for w in ws[:N_WORDS]]
            break
        if isinstance(tw, dict):
            for k in range(1, 10):
                for cand in (f'T{k}', str(k), str(k - 1)):
                    if cand in tw:
                        ws = tw[cand]
                        if isinstance(ws, list):
                            words_per_topic[k - 1] = [str(w) for w in ws[:N_WORDS]]
                        break
            break

    # If still empty, scan for any per-topic word arrays.
    if not any(words_per_topic):
        print("  [warn] Could not locate top-word arrays in nlp_results.json — "
              "edit CANONICAL_NAMES or check the JSON schema")

    return list(zip(names, words_per_topic))


def count_manual_corrections():
    """Count MANUAL_CORRECTIONS entries in src/15_entity_classify.py."""
    path = ROOT / 'src' / '15_entity_classify.py'
    try:
        text = path.read_text()
    except FileNotFoundError:
        return None
    m = re.search(r'MANUAL_CORRECTIONS\s*=\s*\{', text)
    if not m:
        return None
    # Balance braces until matching close.
    start = m.end() - 1
    depth = 0
    for i, ch in enumerate(text[start:], start):
        if ch == '{':
            depth += 1
        elif ch == '}':
            depth -= 1
            if depth == 0:
                body = text[start + 1:i]
                # Count top-level quoted keys (lines with 'foo':).
                return len(re.findall(r"^\s*['\"][^'\"]+['\"]\s*:", body, re.M))
    return None


def entity_network_counts():
    path = ROOT / 'json' / 'entity_network.json'
    try:
        with open(path) as f:
            en = json.load(f)
    except FileNotFoundError:
        return None
    nodes = en.get('nodes') or []
    edges = en.get('edges') or []
    by_kind = {}
    for n in nodes:
        k = n.get('kind') or n.get('type') or 'unknown'
        by_kind[k] = by_kind.get(k, 0) + 1
    return {
        'n_nodes': len(nodes),
        'n_edges': len(edges),
        'by_kind': by_kind,
    }


# ── Slide edits ───────────────────────────────────────────────────────────────

def edit_slide_4(pres):
    """Corpus headline: 695 in Calibre collection · 541 analysed."""
    s = pres.slides[3]
    # Shape[2] = "695" ; Shape[3] = "Books"
    ok1 = set_shape_text(s.shapes[2], "695 · 541")
    ok2 = set_shape_text(s.shapes[3], "Collection · Analysed")
    # Add the OCR-exclusion note to the provisional caveat footer (shape[20]).
    addendum = ("  541 analysed — [2133] *Cybernation and Social Change* excluded "
                "(OCR corruption).")
    if not walk_and_replace(pres, 3, "Findings are provisional — the collection is not exhaustive.",
                            "Findings are provisional — the collection is not exhaustive." + addendum):
        print("  [slide 4] caveat footer not found — addendum skipped")
    return ok1 and ok2


def edit_slide_8(pres, n_scripts, n_methodology_lines, n_decisions_lines):
    """Timeline: update title + extend the last card + refresh footer stats."""
    # Title
    walk_and_replace(pres, 7, "Building CyberneticsNLP — 5 Sessions",
                             "Building CyberneticsNLP — April Sessions (v0.1 → v0.5.1)")

    # Last card currently ends at Apr 14 / v0.4.3. Extend it to cover through v0.5.1.
    s = pres.slides[7]
    # Shape[30]="Apr 14"  [32]="v0.4.3"  [36] = body text
    set_shape_text(s.shapes[30], "Apr 14 → 23")
    set_shape_text(s.shapes[32], "v0.4.3 → v0.5.1")
    set_shape_text(s.shapes[36],
        "Full-text LDA · k=9 canonical\n"
        "Entity network release prep\n"
        "Google Forms survey infra\n"
        "Noise + security cleanup")

    # Footer stats
    # Original: "~34 pipeline scripts  ·  2,000+ line methodology doc  ·
    #            1,600+ line decision log  ·  15 regression tests"
    new_footer = (f"~{n_scripts} pipeline scripts  ·  "
                  f"{n_methodology_lines}+ line methodology doc  ·  "
                  f"{n_decisions_lines}+ line decision log  ·  "
                  f"regression test suite")
    # Match prefix to locate the footer regardless of exact number.
    for shape in s.shapes:
        if not shape.has_text_frame:
            continue
        if "pipeline scripts" in shape.text_frame.text:
            set_shape_text(shape, new_footer)
            break
    return True


def edit_slide_10(pres):
    """Change '695 books' → '541 books' in Abstractive Summaries block."""
    return walk_and_replace(pres, 9, "695 books", "541 books")


def edit_slide_11(pres):
    """Update mean stability and stable-count figures (appear twice)."""
    ok1 = walk_and_replace(pres, 10, "mean stability 0.327, 5/9 topics stable",
                                     "mean stability 0.357, 9/9 topics stable")
    ok2 = walk_and_replace(pres, 10, "5-seed stability complete (mean=0.327)",
                                     "5-seed stability complete (mean=0.357, 9/9 stable)")
    return ok1 or ok2


def edit_slide_13(pres, topic_data):
    """Regenerate the 7 visible topic rows + update the T8/T9 caption."""
    s = pres.slides[12]
    # Layout from inspection:
    #   Row i (i in 0..6) uses shapes [5+6i, 6+6i, 7+6i]
    #   → (T-label, name, words)
    # Check: rows 0..6 have base_shape_idx = [5, 11, 17, 23, 29, 35, 41]
    row_bases = [5, 11, 17, 23, 29, 35, 41]
    for i, base in enumerate(row_bases):
        name, words = topic_data[i]
        # T{i+1} stays as-is but re-set it defensively to guard against stale labels.
        set_shape_text(s.shapes[base], f"T{i+1}")
        set_shape_text(s.shapes[base + 1], name)
        if words:
            set_shape_text(s.shapes[base + 2], " · ".join(words))
        # if words empty, leave the existing top-word string alone

    # Caption line (shape[1]): mentions T8/T9
    t8_name, _ = topic_data[7]
    t9_name, _ = topic_data[8]
    new_caption = (
        f"9-topic solution — 7 shown above. Also: T8 {t8_name} · T9 {t9_name}. "
        f"Names: API-proposed, corrected by Paul; provisional pending multi-rater validation."
    )
    set_shape_text(s.shapes[1], new_caption)
    return True


def edit_slide_16(pres, n_corrections, net):
    """Update the manual-overrides line and inject a network-stats bullet."""
    s = pres.slides[15]
    overrides_text = (
        f"·  {n_corrections}+ post-audit manual corrections across five review batches — "
        f"e.g. 'stelarc' → person, 'digital' → concept, trailing-fragment suppression"
        if n_corrections else
        "·  Hundreds of post-audit manual corrections across five review batches — "
        "e.g. 'stelarc' → person, 'digital' → concept, trailing-fragment suppression"
    )
    # Shape [19] holds the "121 post-audit manual overrides" line.
    set_shape_text(s.shapes[19], overrides_text)

    # If we have live network stats, fold them into the "How Edges…" intro line.
    # Idempotent: strip any prior "Network: … edges" line before prepending a
    # fresh one, so re-runs don't stack multiple stats lines.
    if net:
        kinds = net['by_kind']
        breakdown = f"persons={kinds.get('person', 0)} · concepts={kinds.get('concept', 0)} · " \
                    f"orgs={kinds.get('organisation', kinds.get('org', 0))} · " \
                    f"locations={kinds.get('location', 0)}"
        stats_line = (f"Network: {net['n_nodes']} nodes ({breakdown}) · "
                      f"{net['n_edges']} edges")
        existing = s.shapes[16].text_frame.text
        cleaned = re.sub(r'^Network:.*edges\s*\n?', '', existing, flags=re.M)
        set_shape_text(s.shapes[16], stats_line + "\n" + cleaned.lstrip())
    return True


def edit_slide_5(pres):
    """Shift the 5 question cards down 0.8in so they no longer cover the intro
    paragraph. Shapes [1] and [2] (the intro) stay put. Shapes [3]..[35] move.

    Idempotent: guarded by checking card-1 y position (original ~1.26in;
    post-shift ~2.06in). Skips if already shifted.

    Side-effect on first run: card 1 loses its full-width background stripe
    (shape[1] was doing double duty as both intro text and card 1 bg).
    """
    from pptx.util import Inches, Emu
    s = pres.slides[4]
    # Guard: if card 1 (shape[3]) is already below 1.8in, assume shift applied.
    if Emu(s.shapes[3].top).inches > 1.8:
        return False
    shift = Inches(0.80)
    for i in range(3, 36):
        if i < len(s.shapes):
            sh = s.shapes[i]
            sh.top = sh.top + shift
    return True


def edit_slide_14(pres):
    """Remap 'Dominant: T#' attributions from the defunct 18 April taxonomy to
    the Run C names that actually live in json/nlp_results.json.

    Per-era mappings chosen by fit between top-words and historical narrative:
      Foundational (1940s–60s): T6 Formal Foundations · T8 Cybernetic Management
      Second-order (1960s–70s): T2 Cybernetics and Circularity (bateson) · T3 Biological Systems
      Fragmentation (1980s–90s): T8 Cybernetic Management · T6 Formal Foundations
      Expansion (2000s–25): T5 Cultural Applications · T1 Cybernetics of Political Economy
    """
    remaps = [
        ("Dominant: T6 Formal Foundations · T4 Cybernetic Management",
         "Dominant: T6 Formal Foundations · T8 Cybernetic Management Theory"),
        ("Dominant: T2 Cybernetics of Psychology · T5 Biological Systems",
         "Dominant: T2 Cybernetics and Circularity · T3 Biological Systems Cybernetics"),
        ("Dominant: T4 Cybernetic Management · T6 Formal Foundations",
         "Dominant: T8 Cybernetic Management Theory · T6 Formal Foundations"),
        ("Dominant: T8 Cybernetics of Posthumanism · T9 Cultural Applications",
         "Dominant: T5 Cultural Applications of Cybernetics · T1 Cybernetics of Political Economy"),
    ]
    any_hit = False
    for old, new in remaps:
        if walk_and_replace(pres, 13, old, new):
            any_hit = True
    return any_hit


def edit_slide_22(pres):
    return walk_and_replace(pres, 21, "4 random seeds", "5 random seeds")


def edit_slide_23(pres):
    """Replace the 'Immediately' column items with current sprint work."""
    s = pres.slides[22]
    # Shapes [4..7] are the four arrow-lines.
    new_items = [
        "→  Multi-rater topic naming survey (Google Forms live, ≥2 raters)",
        "→  Cross-run topic stability: graduate compare_topic_runs.py",
        "→  Full spaCy + Wikidata entity pass (remaining)",
        "→  Document unit decision: proceedings/anthology (moratorium — pending signal inventory)",
    ]
    for idx, text in zip([4, 5, 6, 7], new_items):
        set_shape_text(s.shapes[idx], text)
    return True


# ── Main ──────────────────────────────────────────────────────────────────────

def main():
    if not SRC.exists():
        sys.exit(f"Source deck not found: {SRC}")

    print(f"Loading {SRC.name} …")
    pres = Presentation(str(SRC))

    print("Loading live facts …")
    topic_data = load_topic_data()
    n_corr = count_manual_corrections()
    net = entity_network_counts()

    # Footer stat inputs — derive or fall back.
    src_dir = ROOT / 'src'
    n_scripts = len(list(src_dir.glob('*.py')) + list(src_dir.glob('*.sh')))
    try:
        n_methodology_lines = sum(1 for _ in open(ROOT / 'docs' / 'methodology.md'))
    except FileNotFoundError:
        n_methodology_lines = 2000
    try:
        n_decisions_lines = sum(1 for _ in open(ROOT / 'docs' / 'decisions.md'))
    except FileNotFoundError:
        n_decisions_lines = 1600
    # Round down to the nearest hundred for a stable "2,100+" style figure.
    n_methodology_round = (n_methodology_lines // 100) * 100
    n_decisions_round = (n_decisions_lines // 100) * 100

    print(f"  topic_names: {[t[0] for t in topic_data]}")
    print(f"  manual corrections: {n_corr}")
    print(f"  entity network: {net}")
    print(f"  src scripts: {n_scripts}")
    print(f"  methodology lines: {n_methodology_lines} → '{n_methodology_round}+'")
    print(f"  decisions lines:   {n_decisions_lines} → '{n_decisions_round}+'")

    print("\nApplying edits …")
    edits = [
        ("Slide 4  corpus headline",   lambda: edit_slide_4(pres)),
        ("Slide 5  shift cards down",  lambda: edit_slide_5(pres)),
        ("Slide 8  timeline + footer", lambda: edit_slide_8(pres, n_scripts,
                                                             n_methodology_round,
                                                             n_decisions_round)),
        ("Slide 10 summaries count",   lambda: edit_slide_10(pres)),
        ("Slide 11 stability figures", lambda: edit_slide_11(pres)),
        ("Slide 13 topic regen",       lambda: edit_slide_13(pres, topic_data)),
        ("Slide 14 era remap",         lambda: edit_slide_14(pres)),
        ("Slide 16 entity overrides",  lambda: edit_slide_16(pres, n_corr, net)),
        ("Slide 22 seed count",        lambda: edit_slide_22(pres)),
        ("Slide 23 immediately col",   lambda: edit_slide_23(pres)),
    ]
    for label, fn in edits:
        ok = False
        try:
            ok = fn()
        except Exception as exc:
            print(f"  [fail] {label}: {exc}")
            continue
        print(f"  [{'ok' if ok else 'skip'}] {label}")

    print(f"\nWriting {DST.name} …")
    pres.save(str(DST))
    if SRC == DST:
        print(f"Done. Re-edited {DST.name} in place; diff against v2 for a full review.")
    else:
        print(f"Done. Diff {DST.name} against {SRC.name} in PowerPoint to review.")


if __name__ == '__main__':
    main()
